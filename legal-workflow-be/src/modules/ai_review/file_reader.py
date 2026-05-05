"""Extract text content from uploaded files (PDF, DOCX, PPTX, XLSX, TXT)."""

import os


def read_file_content(file_path: str, max_chars: int = 4000) -> str:
    """Read text content from a file. Returns extracted text or empty string."""
    if not os.path.exists(file_path):
        return ""

    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".txt" or ext == ".md" or ext == ".csv":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()[:max_chars]

        elif ext == ".pdf":
            return _read_pdf(file_path, max_chars)

        elif ext == ".docx":
            return _read_docx(file_path, max_chars)

        elif ext == ".pptx":
            return _read_pptx(file_path, max_chars)

        elif ext == ".xlsx":
            return _read_xlsx(file_path, max_chars)

        else:
            return f"[File type {ext} - cannot extract text]"

    except Exception as e:
        return f"[Error reading file: {str(e)[:100]}]"


def _read_pdf(path: str, max_chars: int) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
        if len(text) > max_chars:
            break
    return text[:max_chars]


def _read_docx(path: str, max_chars: int) -> str:
    from docx import Document
    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs)
    return text[:max_chars]


def _read_pptx(path: str, max_chars: int) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        text += f"\n--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text + "\n"
        if len(text) > max_chars:
            break
    return text[:max_chars]


def _read_xlsx(path: str, max_chars: int) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text += f"\n--- Sheet: {sheet} ---\n"
        for row in ws.iter_rows(max_row=50, values_only=True):
            vals = [str(c) if c is not None else "" for c in row]
            text += " | ".join(vals) + "\n"
            if len(text) > max_chars:
                break
    return text[:max_chars]
